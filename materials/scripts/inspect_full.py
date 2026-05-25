"""Full structural inspection of the deck before the rewrite."""
import io, sys
from pptx import Presentation
from pptx.util import Emu

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

p = Presentation('project_presentation_final.pptx')
print(f'Slide size: {Emu(p.slide_width).inches:.2f} x {Emu(p.slide_height).inches:.2f} in')
print()
print('=== LAYOUTS ===')
for i, lo in enumerate(p.slide_layouts):
    print(f'  {i}: {lo.name}')
print()

for idx in range(len(p.slides)):
    s = p.slides[idx]
    print(f'--- Slide {idx+1}  (layout: {s.slide_layout.name}) ---')
    for j, sh in enumerate(s.shapes):
        kind = sh.shape_type
        pos = ''
        try:
            pos = (f'  @({Emu(sh.left).inches:.2f},{Emu(sh.top).inches:.2f}) '
                   f'{Emu(sh.width).inches:.2f}x{Emu(sh.height).inches:.2f}')
        except Exception:
            pos = '  @(n/a)'
        txt = ''
        if sh.has_text_frame:
            txt = ' | ' + sh.text_frame.text[:60].replace('\n', ' / ')
        is_title = ' [TITLE]' if sh == s.shapes.title else ''
        print(f'  [{j}] {kind} name={sh.name!r}{is_title}{pos}{txt}')
    print()
