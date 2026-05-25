# -*- coding: utf-8 -*-
"""Append three Phase-3 evaluation slides to project_presentation_final.pptx:

  * baseline vs fine-tuned comparison (comparison_baseline_bert.png)
  * confusion matrix (confusion_matrix.png)
  * FLAN-T5 ROUGE training curves (flan_t5_rouge_curves.png)

Run after apply_changes.py so the deck's core slides are already in place.
The new slides are appended at the end of the deck before the References slide
so the viva flow stays linear: Results -> Phase-3 evidence -> References.
"""
import io
import sys
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

MATERIALS = Path(__file__).resolve().parent.parent
FIGURES_DIR = MATERIALS / "figures"
PPTX = str(MATERIALS / "deliverables" / "project_presentation_final.pptx")
NAVY = RGBColor(0x1A, 0x32, 0x60)
DK = RGBColor(0x3D, 0x3D, 0x3D)


def make_slide(prs: Presentation, title_text: str, image_path: str, caption: str) -> None:
    """Add a blank slide with a title, a centred image and a caption below."""
    blank = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(blank)

    sw, sh = prs.slide_width, prs.slide_height

    # Title
    title = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), sw - Inches(0.8), Inches(0.7))
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title_text
    r.font.bold = True
    r.font.size = Pt(26)
    r.font.color.rgb = NAVY

    # Image — centred, max 6.6 inches wide
    img_w = Inches(8.5)
    img_left = (sw - img_w) // 2
    pic = slide.shapes.add_picture(str(FIGURES_DIR / image_path), img_left, Inches(1.15), width=img_w)
    # Constrain height: if too tall, shrink and re-center horizontally
    if pic.height > Inches(5.4):
        ratio = Inches(5.4) / pic.height
        pic.height = Inches(5.4)
        pic.width = int(pic.width * ratio)
        pic.left = (sw - pic.width) // 2

    # Caption underneath
    cap_top = pic.top + pic.height + Inches(0.10)
    cap = slide.shapes.add_textbox(Inches(0.4), cap_top, sw - Inches(0.8), Inches(0.7))
    ctf = cap.text_frame
    ctf.word_wrap = True
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = caption
    cr.font.size = Pt(12)
    cr.font.italic = True
    cr.font.color.rgb = DK


def main() -> int:
    prs = Presentation(PPTX)
    initial_count = len(prs.slides)

    make_slide(
        prs,
        "Baseline vs Fine-tuned: macro F1 lift",
        "comparison_baseline_bert.png",
        "TF-IDF + Logistic Regression macro F1 = 0.61 vs fine-tuned Legal-BERT "
        "macro F1 = 0.902 on the same held-out 419-clause test set (a 48% relative lift).",
    )
    make_slide(
        prs,
        "Confusion matrix — Legal-BERT held-out test set",
        "confusion_matrix.png",
        "Off-diagonal mass is concentrated on OTHER vs INTELLECTUAL_PROPERTY and on "
        "WARRANTY boundary cases — both are over-prediction patterns, not missed classes.",
    )
    make_slide(
        prs,
        "FLAN-T5 rewriter — ROUGE across training epochs",
        "flan_t5_rouge_curves.png",
        "Validation ROUGE-1 / ROUGE-2 / ROUGE-L for the fine-tuned FLAN-T5-small "
        "clause simplifier (freak3123/flan-t5-simplify-v1). Final: 0.394 / 0.231 / 0.352.",
    )
    make_slide(
        prs,
        "Legal-BERT fine-tuning — losses + macro F1 over 4 epochs",
        "legal_bert_training_curves.png",
        "Training and validation losses fall monotonically; validation macro F1 climbs "
        "from 0.512 at epoch 1 to a best of 0.885 at epoch 3 (the checkpoint that was saved).",
    )
    make_slide(
        prs,
        "Sample plain-English rewrites — FLAN-T5",
        "rewrite_examples.png",
        "Three verbatim model outputs on representative clauses. The rewriter preserves "
        "core terms (30-day notice, royalty-free licence, third-party sharing) while "
        "stripping out legalese.",
    )

    prs.save(PPTX)
    print(f"appended 3 slides ({initial_count} -> {len(prs.slides)})")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
