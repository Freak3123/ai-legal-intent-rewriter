"""Inspect theme fonts/colours of the template and the current deck,
and the body-text formatting of untouched slides (the template look)."""
import io, sys, zipfile, re
from pptx import Presentation
from pptx.util import Pt

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')


def theme_info(path, label):
    print(f'===== {label}: {path} =====')
    with zipfile.ZipFile(path) as z:
        themes = [n for n in z.namelist() if re.match(r'ppt/theme/theme\d+\.xml', n)]
        for tn in themes[:1]:
            xml = z.read(tn).decode('utf-8', 'replace')
            mj = re.search(r'<a:majorFont>.*?<a:latin typeface="([^"]*)"', xml, re.S)
            mn = re.search(r'<a:minorFont>.*?<a:latin typeface="([^"]*)"', xml, re.S)
            print(f'  major font: {mj.group(1) if mj else "?"}')
            print(f'  minor font: {mn.group(1) if mn else "?"}')
            cols = re.findall(r'<a:(dk1|lt1|dk2|lt2|accent[1-6])>\s*<a:(srgbClr|sysClr)[^>]*?(?:val|lastClr)="([0-9A-Fa-f]{6})"', xml)
            for name, _, val in cols:
                print(f'  {name}: #{val}')
    print()


for path, label in [('End Term Project Presentation Template.pptx', 'TEMPLATE'),
                     ('project_presentation_final.preedit.pptx', 'DECK (pristine)')]:
    try:
        theme_info(path, label)
    except Exception as e:
        print(f'  ERROR {label}: {e}\n')

print('===== body-text formatting of untouched deck slides =====')
p = Presentation('project_presentation_final.preedit.pptx')
for idx in (2, 4):  # slide 3 Introduction, slide 5 Motivations (0-based)
    s = p.slides[idx]
    print(f'--- Slide {idx + 1}: {s.shapes.title.text if s.shapes.title else "?"} ---')
    for sh in s.shapes:
        if sh.name.startswith('Content Placeholder') and sh.has_text_frame:
            for pi, para in enumerate(sh.text_frame.paragraphs[:3]):
                for r in para.runs:
                    sz = r.font.size.pt if r.font.size else None
                    col = None
                    try:
                        col = r.font.color.rgb
                    except Exception:
                        col = '(theme/inherit)'
                    print(f'  para{pi} run: text={r.text[:30]!r} font={r.font.name} '
                          f'size={sz} bold={r.font.bold} color={col}')
        if sh == s.shapes.title:
            for r in sh.text_frame.paragraphs[0].runs:
                sz = r.font.size.pt if r.font.size else None
                print(f'  TITLE run: font={r.font.name} size={sz} bold={r.font.bold}')
print()
print('===== layout content-placeholder list sizes =====')
for li, layout in enumerate(p.slide_layouts):
    if layout.name in ('Title and Content',):
        for ph in layout.placeholders:
            if ph.placeholder_format.idx == 1 or 'Content' in ph.name:
                txt = ph.text_frame
                for para in txt.paragraphs[:3]:
                    print(f'  layout "{layout.name}" lvl{para.level}: '
                          f'size={para.font.size.pt if para.font.size else None}')
