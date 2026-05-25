"""Remove the previously-inserted chart on the Results slide and re-embed the
new wider/shorter results_chart.png so it fits inside the slide."""
from pptx import Presentation
from pptx.util import Inches

prs = Presentation('project_presentation_final.pptx')
slide_w = prs.slide_width

# Find the Results slide by title.
target = None
for i, s in enumerate(prs.slides):
    title_shape = s.shapes.title
    if title_shape and title_shape.text.strip().lower() == 'results':
        target = (i, s)
        break

assert target, 'Results slide not found'
i, slide = target
print(f'Results slide is at index {i+1}')

# Drop any existing pictures on the slide.
removed = 0
for shape in list(slide.shapes):
    if shape.shape_type == 13:  # PICTURE
        sp = shape._element
        sp.getparent().remove(sp)
        removed += 1
print(f'Removed {removed} existing picture(s)')

# Insert the refreshed chart, centred.
img_w = Inches(12.5)
img_left = (slide_w - img_w) // 2
img_top = Inches(1.95)
slide.shapes.add_picture('results_chart.png', img_left, img_top, width=img_w)

prs.save('project_presentation_final.pptx')
print('Saved.')
