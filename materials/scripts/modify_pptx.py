"""Modify project_presentation_final.pptx:

1. Replace Slide 7 (Literature Review) and insert additional slides so that
   all 35 references are listed in the same 6-column schema as the original
   placeholder table, split across slides (7 rows per slide).
2. Replace the textual bullets on Slide 17 (Results) with the generated
   results_chart.png embed plus a short caption.

The pptx is edited in place. A .backup.pptx was already taken alongside.
"""
import copy
import io
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from lit_review_rows import ROWS

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

PPTX_PATH = 'project_presentation_final.pptx'
CHART_PATH = 'results_chart.png'

HEADER_COLS = ['Sl. No.', 'Author & Year', 'Methods / Techniques',
               'Dataset', 'Research Outcomes', 'Findings']

ROWS_PER_SLIDE = 7
CHUNK_COUNT = (len(ROWS) + ROWS_PER_SLIDE - 1) // ROWS_PER_SLIDE  # = 5

prs = Presentation(PPTX_PATH)
slide_w, slide_h = prs.slide_width, prs.slide_height


# ---------- helpers ----------------------------------------------------------

def clone_slide_after(prs, source_slide, insert_after_idx):
    """Duplicate source_slide and move the duplicate so it sits at index
    insert_after_idx + 1 in the slide order. Returns the new slide."""
    # Add at end using the same layout.
    new = prs.slides.add_slide(source_slide.slide_layout)
    # Copy every shape's XML.
    from copy import deepcopy
    # Remove default shapes (title/content placeholders) on the new slide first.
    for shape in list(new.shapes):
        sp = shape._element
        sp.getparent().remove(sp)
    for shape in source_slide.shapes:
        new.shapes._spTree.insert_element_before(deepcopy(shape._element), 'p:extLst')
    # Reorder so the new slide is right after source_slide.
    sldIdLst = prs.slides._sldIdLst
    slides = list(sldIdLst)
    # The newly added slide id is the last entry.
    new_id = slides[-1]
    sldIdLst.remove(new_id)
    # Insert at insert_after_idx + 1.
    sldIdLst.insert(insert_after_idx + 1, new_id)
    return new


def set_cell(cell, text, *, bold=False, font_size=10, bg=None, fg=None,
             align=PP_ALIGN.LEFT):
    cell.text = ''  # clear default paragraph contents
    tf = cell.text_frame
    tf.word_wrap = True
    # margins
    cell.margin_left = Inches(0.05)
    cell.margin_right = Inches(0.05)
    cell.margin_top = Inches(0.03)
    cell.margin_bottom = Inches(0.03)
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if fg is not None:
        run.font.color.rgb = fg
    if bg is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


def build_lit_review_table(slide, rows, slide_label):
    """Place a 6-col table with header + rows on the given slide."""
    # Remove the existing content placeholder body text and any existing table.
    for shape in list(slide.shapes):
        if shape.name in ('Content Placeholder 2', 'Content Placeholder 6'):
            tf = shape.text_frame
            tf.clear()
            # Put a small caption above the table area.
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = f'Table 1: Summary of existing approaches and research findings  —  {slide_label}'
            r.font.size = Pt(13)
            r.font.bold = True
        elif shape.has_table:
            sp = shape._element
            sp.getparent().remove(sp)

    n_rows = len(rows) + 1  # +1 for header
    n_cols = 6
    left = Inches(0.40)
    top = Inches(1.65)
    width = Inches(12.55)
    height = Inches(5.10)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # Column widths (sum should ~= 12.55").
    col_widths_in = [0.55, 1.55, 2.40, 2.10, 2.95, 3.00]
    for i, w in enumerate(col_widths_in):
        table.columns[i].width = Inches(w)

    # Row heights — header smaller, body rows even.
    header_h = Inches(0.45)
    body_h = Inches((5.10 - 0.45) / len(rows))
    table.rows[0].height = header_h
    for r in range(1, n_rows):
        table.rows[r].height = body_h

    # Header.
    header_bg = RGBColor(0x1F, 0x3A, 0x5F)
    header_fg = RGBColor(0xFF, 0xFF, 0xFF)
    for i, col_name in enumerate(HEADER_COLS):
        set_cell(table.cell(0, i), col_name, bold=True, font_size=11,
                 bg=header_bg, fg=header_fg, align=PP_ALIGN.CENTER)

    # Body.
    alt_bg = RGBColor(0xF2, 0xF5, 0xF9)
    for ri, row in enumerate(rows, start=1):
        zebra = alt_bg if ri % 2 == 0 else None
        for ci, val in enumerate(row):
            align = PP_ALIGN.CENTER if ci == 0 else PP_ALIGN.LEFT
            bold = ci == 0
            set_cell(table.cell(ri, ci), val, font_size=9, bold=bold,
                     bg=zebra, align=align)


def renumber_slide_numbers(prs):
    """Update any 'Slide Number Placeholder' text shapes so the numbers match
    the new ordinal positions."""
    for idx, slide in enumerate(prs.slides):
        target = idx + 1
        for shape in slide.shapes:
            if shape.name.startswith('Slide Number Placeholder') and shape.has_text_frame:
                tf = shape.text_frame
                # Some templates use a field; we just rewrite the visible text.
                if tf.paragraphs and tf.paragraphs[0].runs:
                    for run in tf.paragraphs[0].runs:
                        run.text = str(target)
                else:
                    tf.text = str(target)


# ---------- step 1: literature review (slide index 6 → slide 7) -------------

lit_slide_idx = 6  # 0-based
lit_slide = prs.slides[lit_slide_idx]

# Chunk rows into groups of ROWS_PER_SLIDE.
chunks = [ROWS[i:i + ROWS_PER_SLIDE] for i in range(0, len(ROWS), ROWS_PER_SLIDE)]
print(f'Lit review chunks: {[len(c) for c in chunks]} (total {sum(len(c) for c in chunks)})')

# First chunk replaces the existing slide 7.
build_lit_review_table(lit_slide, chunks[0],
                       f'part 1 of {CHUNK_COUNT}  —  references 1-{len(chunks[0])}')

# Clone the slide for each remaining chunk, inserting immediately after the previous.
prev_idx = lit_slide_idx
prev_slide = lit_slide
running = len(chunks[0])
for ci, chunk in enumerate(chunks[1:], start=2):
    new_slide = clone_slide_after(prs, prev_slide, prev_idx)
    label = f'part {ci} of {CHUNK_COUNT}  —  references {running + 1}-{running + len(chunk)}'
    build_lit_review_table(new_slide, chunk, label)
    prev_idx += 1
    prev_slide = new_slide
    running += len(chunk)

# ---------- step 2: results slide -------------------------------------------

# After inserting (CHUNK_COUNT - 1) new slides, original slide 17 has shifted.
new_results_idx = 16 + (CHUNK_COUNT - 1)  # 16 + 4 = 20
results_slide = prs.slides[new_results_idx]
assert results_slide.shapes.title is not None
print(f'Results slide is now at index {new_results_idx + 1} '
      f"(title='{results_slide.shapes.title.text}')")

# Clear the body placeholder text and any existing pictures, then insert the chart.
chart_placeholder = None
for shape in list(results_slide.shapes):
    if shape.name in ('Content Placeholder 6', 'Content Placeholder 2'):
        chart_placeholder = shape
        # Empty its text and replace with a short caption line.
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = ('Per-class classifier performance under stratified cross-validation '
                  '(measured; source: research_paper_final.docx, Table 9).')
        r.font.size = Pt(12)
        r.font.bold = True

# Insert chart image, centred under the caption.
img_w = Inches(11.6)
img_left = (slide_w - img_w) // 2
img_top = Inches(1.90)
results_slide.shapes.add_picture(CHART_PATH, img_left, img_top, width=img_w)

# ---------- step 3: renumber slide numbers ----------------------------------

renumber_slide_numbers(prs)

prs.save(PPTX_PATH)
print(f'Saved {PPTX_PATH}')
print(f'Total slides now: {len(prs.slides)}')
