from pptx import Presentation

p = Presentation('project_presentation_final.pptx')
print('Total slides:', len(p.slides))
print()

for i, slide in enumerate(p.slides):
    print(f'===== Slide {i+1} =====')
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text
                if text.strip():
                    print(text)
        elif shape.shape_type == 19:  # Table
            print('[TABLE]')
            tbl = shape.table
            for row in tbl.rows:
                row_text = ' | '.join(cell.text for cell in row.cells)
                print(row_text)
    print()
