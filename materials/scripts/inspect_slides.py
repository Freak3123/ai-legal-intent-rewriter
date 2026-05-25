"""Inspect existing slide layouts and shape positions for the slides we will edit."""
from pptx import Presentation
from pptx.util import Emu
import sys, io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

p = Presentation('project_presentation_final.pptx')

print(f"Slide dimensions: {p.slide_width} x {p.slide_height} EMU = {Emu(p.slide_width).inches:.2f} x {Emu(p.slide_height).inches:.2f} inches")
print()

for slide_idx in [6, 16, 18]:  # Slide 7, 17, 19 (0-indexed)
    slide = p.slides[slide_idx]
    print(f"=== Slide {slide_idx+1} ===")
    print(f"Layout: {slide.slide_layout.name}")
    for j, shape in enumerate(slide.shapes):
        print(f"  Shape {j}: {shape.shape_type} name='{shape.name}'")
        print(f"    pos=({Emu(shape.left).inches:.2f}\", {Emu(shape.top).inches:.2f}\") size=({Emu(shape.width).inches:.2f}\" x {Emu(shape.height).inches:.2f}\")")
        if shape.has_text_frame:
            txt = shape.text_frame.text[:100].replace('\n', ' | ')
            print(f"    text: {txt}")
    print()

# Also list slide layouts
print("=== Available layouts ===")
for i, layout in enumerate(p.slide_layouts):
    print(f"  Layout {i}: {layout.name}")
