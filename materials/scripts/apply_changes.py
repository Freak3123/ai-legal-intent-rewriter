# -*- coding: utf-8 -*-
"""Apply the requested presentation rewrite to project_presentation_final.pptx.

Re-runnable: a pristine snapshot (project_presentation_final.preedit.pptx) is
created on first run; every run restores from it.
  python apply_changes.py            -> writes the real deck (PowerPoint must be CLOSED)
  python apply_changes.py out.pptx   -> writes out.pptx (dry run; PowerPoint may be open)

Text uses the deck theme: Gill Sans MT (inherited) + theme navy #1A3260.
Body sizes are role-uniform: 18 standard / 13 dense / 12 diagram columns / 10.5 refs.
"""
import io
import os
import shutil
import sys
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from PIL import Image

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pathlib import Path
MATERIALS = Path(__file__).resolve().parent.parent
FIGURES_DIR = MATERIALS / 'figures'
PPTX    = str(MATERIALS / 'deliverables' / 'project_presentation_final.pptx')
PREEDIT = str(MATERIALS / 'backups' / 'project_presentation_final.preedit.pptx')
# Make bare filename references to PNGs (e.g. 'results_chart.png') resolve
# correctly when passed to add_picture by chdir-ing to the figures folder.
os.chdir(FIGURES_DIR)

NAVY = RGBColor(0x1A, 0x32, 0x60)   # theme accent1 — used for emphasis only
DK   = RGBColor(0x3D, 0x3D, 0x3D)   # theme dk2 — secondary text

# role-uniform body sizes
SZ_STD, SZ_DENSE, SZ_COL, SZ_REF = 18, 13, 12, 10.5

OUT = sys.argv[1] if len(sys.argv) > 1 else PPTX
if not os.path.exists(PREEDIT):
    shutil.copyfile(PPTX, PREEDIT)
    print(f'created pristine snapshot: {PREEDIT}')
prs = Presentation(PREEDIT)
SLIDE_W, SLIDE_H = prs.slide_width, prs.slide_height


# --------------------------------------------------------------------------
# helpers
# --------------------------------------------------------------------------
def slide_by_title(text):
    for s in prs.slides:
        t = s.shapes.title
        if t is not None and text.lower() in (t.text or '').lower():
            return s
    raise KeyError(text)


def body_placeholder(slide):
    for sh in slide.shapes:
        if sh.name.startswith('Content Placeholder'):
            return sh
    raise KeyError('body placeholder')


def shape_by_name(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    raise KeyError(name)


def _no_bullet(p):
    pPr = p._p.get_or_add_pPr()
    for tag in ('a:buChar', 'a:buAutoNum', 'a:buNone'):
        for e in pPr.findall(qn(tag)):
            pPr.remove(e)
    pPr.append(pPr.makeelement(qn('a:buNone'), {}))


def _hanging(p, mar=0.30):
    pPr = p._p.get_or_add_pPr()
    pPr.set('marL', str(int(Inches(mar))))
    pPr.set('indent', str(-int(Inches(mar))))


def fill(tf, specs, *, clear=True, word_wrap=True):
    """specs: list of dicts. Each: text | runs (list of (text,bold[,color])),
       size, bold, italic, color, align, space_before, space_after,
       bullet(False to suppress), hanging(bool).
       A run with no colour inherits the theme text colour (template-matching)."""
    if clear:
        tf.clear()
    tf.word_wrap = word_wrap
    for i, spec in enumerate(specs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        runs = spec.get('runs') or [(spec.get('text', ''), spec.get('bold', False))]
        for rspec in runs:
            rt = rspec[0]
            rb = rspec[1] if len(rspec) > 1 else False
            rc = rspec[2] if len(rspec) > 2 else spec.get('color')
            r = p.add_run()
            r.text = rt
            r.font.size = Pt(spec.get('size', SZ_STD))
            r.font.bold = rb
            if spec.get('italic'):
                r.font.italic = True
            if rc is not None:
                r.font.color.rgb = rc
        if 'align' in spec:
            p.alignment = spec['align']
        if spec.get('space_before') is not None:
            p.space_before = Pt(spec['space_before'])
        if spec.get('space_after') is not None:
            p.space_after = Pt(spec['space_after'])
        if spec.get('bullet') is False:
            _no_bullet(p)
        if spec.get('hanging'):
            _hanging(p)


def add_picture(slide, path, *, left, top, width):
    w, h = Image.open(path).size
    pic_w = Inches(width)
    pic_h = int(pic_w * h / w)
    return slide.shapes.add_picture(path, Inches(left), Inches(top),
                                    width=pic_w, height=pic_h)


def drop_pictures(slide):
    for sh in list(slide.shapes):
        if sh.shape_type == 13:  # PICTURE
            sh._element.getparent().remove(sh._element)


def clone_slide_after(source_slide, insert_after_idx):
    new = prs.slides.add_slide(source_slide.slide_layout)
    for sh in list(new.shapes):
        sh._element.getparent().remove(sh._element)
    for sh in source_slide.shapes:
        new.shapes._spTree.insert_element_before(deepcopy(sh._element), 'p:extLst')
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    new_id = ids[-1]
    sldIdLst.remove(new_id)
    sldIdLst.insert(insert_after_idx + 1, new_id)
    return new


def renumber():
    for idx, slide in enumerate(prs.slides):
        for sh in slide.shapes:
            if sh.name.startswith('Slide Number Placeholder') and sh.has_text_frame:
                tf = sh.text_frame
                if tf.paragraphs and tf.paragraphs[0].runs:
                    for run in tf.paragraphs[0].runs:
                        run.text = ''
                    tf.paragraphs[0].runs[0].text = str(idx + 1)
                else:
                    tf.text = str(idx + 1)


def bullets(items, size, *, lead_navy=False, char='•  ', space_after=7):
    """Build literal-bullet textbox paragraphs. items: str OR (lead, rest)."""
    out = []
    for it in items:
        if isinstance(it, tuple):
            runs = [(char, True, NAVY), (it[0], True, NAVY), (it[1], False)]
        else:
            runs = [(char, True, NAVY), (it, False)]
        out.append(dict(size=size, space_after=space_after, bullet=False, runs=runs))
    return out


# --------------------------------------------------------------------------
# 1. global text replacements ("Final" -> "Capstone", typo fix)
# --------------------------------------------------------------------------
REPLACE = {
    'Final Year Research Project': 'Capstone Year Research Project',
    'Final-Year Research Project': 'Capstone Year Research Project',
    'final-year': 'capstone-year',
    'Final-year': 'Capstone-year',
    'Conclustion': 'Conclusion',
}


def _fix_tf(tf):
    for p in tf.paragraphs:
        for r in p.runs:
            t = r.text
            for a, b in REPLACE.items():
                if a in t:
                    t = t.replace(a, b)
            if t != r.text:
                r.text = t


for s in prs.slides:
    for sh in s.shapes:
        if sh.has_text_frame:
            _fix_tf(sh.text_frame)
        if sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    _fix_tf(cell.text_frame)
print('1. global "Final" -> "Capstone" replacement done')


# --------------------------------------------------------------------------
# 2. Slide 1 - project title
# --------------------------------------------------------------------------
s1 = prs.slides[0]
tb = shape_by_name(s1, 'TextBox 8')
tb.top, tb.height = Inches(1.02), Inches(2.30)
fill(tb.text_frame, [
    dict(text='END-TERM PRESENTATION', size=14, bold=True, color=DK,
         align=PP_ALIGN.CENTER, space_after=6),
    dict(text='AI Legal Intent Rewriter', size=32, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER, space_after=4),
    dict(text='with a Multi-Stage NLP Pipeline for Plain-Language '
              'Interpretation and Risk Flagging of Legal Documents',
         size=17, color=NAVY, align=PP_ALIGN.CENTER),
])
print('2. slide 1 project title set')


# --------------------------------------------------------------------------
# 3. Slide 4 - Problem Statement (SMART, single paragraph)
# --------------------------------------------------------------------------
s4 = body_placeholder(slide_by_title('Problem Statement'))
fill(s4.text_frame, [dict(size=SZ_STD, align=PP_ALIGN.JUSTIFY, bullet=False, runs=[
    ("Legal documents — tenancy agreements, employment contracts and online "
     "terms of service — are written in dense ‘legalese’ that non-experts "
     "cannot interpret, so most people enter binding agreements without genuine "
     "informed consent. ", False),
    ("Specifically", True),
    (", this project builds a multi-stage NLP pipeline that, given a digital-PDF "
     "or plain-text contract, segments it into clauses and for each clause assigns "
     "one of twelve legal categories, extracts key entities, and produces a "
     "plain-language rewrite and a transparent risk level. Progress is ", False),
    ("measurable", True),
    (" against defined targets — ≥85% clause-boundary accuracy, ≥0.80 macro-F1 "
     "for classification and ≥0.75 entity-level F1 for NER — verified by an "
     "automated test suite. The work is ", False),
    ("assignable", True),
    (": it is owned by the four-member Group 3211, with classification and NER, "
     "rewriting and risk flagging, the web dashboard, and dataset and ingestion "
     "each assigned to a named member under one supervisor. It is ", False),
    ("realistic", True),
    (" — scoped to English-language consumer contracts and built entirely on "
     "free, open-source tools and free-tier hosting that run on commodity hardware "
     "with no dedicated GPU. Finally it is ", False),
    ("time-bound", True),
    (", delivered across an 18-week capstone schedule. The integrated pipeline, "
     "the transformer fine-tuning and the end-to-end deployment are all complete; "
     "the held-out macro-F1 of 0.902 exceeds the project target.", False),
])])
print('3. slide 4 problem statement (SMART paragraph) set')


# --------------------------------------------------------------------------
# 4. Slide 6 - Project Objectives (SMART bullets)
# --------------------------------------------------------------------------
s6 = body_placeholder(slide_by_title('Project Objectives'))
OBJ = [
    "Ingest PDF and plain-text legal documents, with OCR for scanned inputs, "
    "producing clean normalised text for the pipeline.",
    "Segment each document into individual clauses with character-accurate "
    "offsets — target at least 85% boundary-detection accuracy.",
    "Classify every clause into one of 12 legal categories — target a macro "
    "F1-score of at least 0.80 on the supported categories.",
    "Extract 6 entity types (dates, amounts, parties, rights, obligations, "
    "conditions) — target at least 0.75 entity-level F1.",
    "Generate faithful plain-language rewrites of each clause, evaluated by "
    "readability and similarity metrics.",
    "Flag risky clauses with a transparent, explainable rule engine "
    "(low, medium or high, every trigger named).",
    "Present every result in a responsive, clause-by-clause web dashboard.",
    "Validate all components with quantitative metrics and an automated test "
    "suite, within the 18-week capstone timeline.",
]
fill(s6.text_frame, [dict(text=t, size=SZ_STD, space_after=6) for t in OBJ])
print('4. slide 6 objectives (SMART bullets) set')


# --------------------------------------------------------------------------
# 5. Slide 12 - Research Gap & Improvements (detailed, with sources)
# --------------------------------------------------------------------------
s12 = body_placeholder(slide_by_title('Research Gap'))
fill(s12.text_frame, [
    dict(size=SZ_DENSE, space_after=3, runs=[
        ("1.  Fragmentation of capability — ", True, NAVY),
        ("existing tools each solve only one sub-task. Classification benchmarks "
         "(CUAD [11], LEDGAR [34], LexGLUE [33]), text-simplification systems "
         "(Cemri et al. [17]; legal-contract simplification [18]) and unfair-clause "
         "detectors (CLAUDETTE — Lippi et al. [21]) are each built and evaluated "
         "in isolation — yet a non-expert reader needs all of them at once.", False),
    ]),
    dict(size=SZ_DENSE, space_after=11, runs=[
        ("Our improvement — ", True, NAVY),
        ("one integrated, clause-level pipeline joining classification, entity "
         "extraction, rewriting and risk flagging through a single JSON contract.", False),
    ]),
    dict(size=SZ_DENSE, space_after=3, runs=[
        ("2.  Summarisation is not comprehension — ", True, NAVY),
        ("most consumer-facing work (e.g. Manor & Li [20]) produces only a shorter "
         "version of a document. Shortening text the reader still cannot interpret "
         "does not, by itself, create understanding.", False),
    ]),
    dict(size=SZ_DENSE, space_after=11, runs=[
        ("Our improvement — ", True, NAVY),
        ("genuine clause-by-clause explanation: each clause is classified, "
         "rewritten in plain language and risk-rated, not merely abbreviated.", False),
    ]),
    dict(size=SZ_DENSE, space_after=3, runs=[
        ("3.  Opaque risk assessment — ", True, NAVY),
        ("several unfair-clause detectors are neural and opaque (Ruas et al. [22]; "
         "B2C abusive-clause detector [23]; chain-of-thought LLM prompting [24]), so "
         "the user cannot see why a clause was flagged.", False),
    ]),
    dict(size=SZ_DENSE, space_after=2, runs=[
        ("Our improvement — ", True, NAVY),
        ("a deterministic three-layer rule-based risk engine in which every "
         "trigger is named and inspectable.", False),
    ]),
])
print('5. slide 12 research gap (detailed + sources) set')


# --------------------------------------------------------------------------
# 6. Slide 13 - Proposed Solution & Architecture (prose, no arrows, + diagram)
# --------------------------------------------------------------------------
s13 = slide_by_title('Proposed Solution')
tb13 = shape_by_name(s13, 'TextBox 10')
tb13.left, tb13.top = Inches(0.29), Inches(1.46)
tb13.width, tb13.height = Inches(5.05), Inches(5.35)
S13 = [
    "The system is two cooperating services that exchange data over a stable, "
    "versioned HTTP / JSON contract (the /v1/analyze and /v1/health endpoints).",
    "The Next.js 15 web application owns everything the user touches — the "
    "dashboard, Auth.js v5 login, document upload, MongoDB persistence and "
    "orchestration.",
    "The FastAPI service (Python 3.11) owns all NLP and machine-learning "
    "inference; the web application never runs a model itself.",
    "For each document the ML service runs six stages in order: it ingests and "
    "normalises the text, then segments it into clauses, and then processes every "
    "clause in turn — classifying it, extracting its entities, rewriting it in "
    "plain language and assigning a risk level.",
    "The structured JSON result is stored in MongoDB Atlas, which holds the "
    "users, documents, analyses and clauses.",
    "Digital-PDF text is extracted in the browser with pdf.js; scanned PDFs are "
    "routed to server-side OCR before analysis.",
    "The split is deliberate: serverless platforms cap the bundle size below the "
    "model size and JavaScript inference is slower, while a self-hosted service "
    "avoids the cold starts and rate limits of hosted APIs.",
]
fill(tb13.text_frame, bullets(S13, SZ_COL, space_after=6))
drop_pictures(s13)
add_picture(s13, 'architecture_diagram.png', left=5.62, top=1.55, width=7.20)
print('6. slide 13 proposed solution (prose, no arrows) + architecture diagram')


# --------------------------------------------------------------------------
# 7. Slide 14 - Dataset Description (sources + download location)
# --------------------------------------------------------------------------
s14 = body_placeholder(slide_by_title('Dataset Description'))
DS = [
    ("Primary source — ", "CUAD (Contract Understanding Atticus Dataset), "
     "Hendrycks et al., NeurIPS 2021, created by The Atticus Project."),
    ("Downloaded from — ", "the official Atticus Project release on GitHub, "
     "github.com/TheAtticusProject/cuad (also mirrored on Hugging Face Datasets "
     "as theatticusproject/cuad-qa)."),
    ("Scale — ", "510 commercial contracts with 13,000+ expert annotations "
     "across 41 fine-grained clause categories."),
    ("Licence — ", "released under Creative Commons CC BY 4.0, free for "
     "academic and commercial use."),
    ("Structure — ", "CUAD frames contract review as question-answering: each "
     "example pairs a category question with an answer span in a contract."),
    ("Taxonomy — ", "the 41 CUAD categories are mapped onto a purpose-built "
     "12-category taxonomy; uncertain mappings are sent to OTHER."),
    ("Seed source — ", "clauses_seed.csv, 72 hand-curated clauses (6 per "
     "category) covering the categories CUAD under-supports."),
    ("Final training set — ", "4,034 labelled clauses — 3,962 from CUAD after "
     "balanced sampling, plus the 72 seed clauses."),
]
fill(s14.text_frame, [
    dict(size=SZ_STD, space_after=6,
         runs=[(lead, True, NAVY), (rest, False)]) for lead, rest in DS
])
print('7. slide 14 dataset description (sources + download URL) set')


# --------------------------------------------------------------------------
# 8. Slide 15 - Data Preprocessing (bullets + flow image)
# --------------------------------------------------------------------------
s15slide = slide_by_title('Data Preprocessing')
s15 = body_placeholder(s15slide)
s15.left, s15.top = Inches(0.29), Inches(1.18)
s15.width, s15.height = Inches(12.84), Inches(1.85)
PRE = [
    "CUAD answer spans are parsed directly from the official release; empty "
    "answers and spans under 20 characters are discarded.",
    "Spans longer than 2,000 characters are truncated; document-metadata "
    "categories (names, dates, parties) are filtered out.",
    "The 41 CUAD categories are mapped to the project's 12-category taxonomy "
    "and the text is normalised (accents, whitespace, casing).",
    "Balanced sampling (at most 800 spans per class) and the 72 merged seed "
    "clauses give the final 4,034-example training set.",
]
fill(s15.text_frame, [dict(text=t, size=SZ_DENSE, space_after=3) for t in PRE])
drop_pictures(s15slide)
add_picture(s15slide, 'preprocessing_flow.png', left=0.37, top=3.18, width=12.60)
print('8. slide 15 preprocessing bullets + flow image set')


# --------------------------------------------------------------------------
# 9. Slide 16 - EDA (bullets + distribution chart)
# --------------------------------------------------------------------------
s16slide = slide_by_title('Exploratory Data Analysis')
s16 = body_placeholder(s16slide)
s16.left, s16.top = Inches(0.29), Inches(1.18)
s16.width, s16.height = Inches(12.84), Inches(1.20)
EDA = [
    "9,118 usable clause spans were extracted from CUAD; balanced sampling "
    "plus 72 seed clauses yields 4,034 labelled examples.",
    "Seven categories are well-supported (181–806 examples each); five "
    "categories rely on just 6 seed examples each — a severe class imbalance "
    "and the key modelling challenge.",
]
fill(s16.text_frame, [dict(text=t, size=SZ_DENSE, space_after=3) for t in EDA])
drop_pictures(s16slide)
add_picture(s16slide, 'eda_label_distribution.png', left=0.97, top=2.45, width=11.40)
print('9. slide 16 EDA bullets + distribution chart set')


# --------------------------------------------------------------------------
# 10. Slide 17 - Model Selection (bullets + comparison image)
# --------------------------------------------------------------------------
s17slide = slide_by_title('Model Selection')
s17 = body_placeholder(s17slide)
s17.left, s17.top = Inches(0.29), Inches(1.16)
s17.width, s17.height = Inches(12.84), Inches(1.05)
MS = [
    "Task — 12-class clause classification, plus a separate clause-rewriting "
    "model; both must run on free CPU-tier hosting at inference time.",
    "Chosen models — fine-tuned nlpaueb/legal-bert-base-uncased classifier "
    "(macro F1 0.902) and fine-tuned google/flan-t5-small rewriter (ROUGE-L "
    "0.352), both pushed to the Hugging Face Hub. A TF-IDF + Logistic "
    "Regression baseline (macro F1 0.61) is retained as a runtime fallback.",
]
fill(s17.text_frame, [dict(text=t, size=SZ_DENSE, space_after=3) for t in MS])
drop_pictures(s17slide)
add_picture(s17slide, 'model_selection.png', left=1.37, top=2.30, width=10.60)
print('10. slide 17 model selection bullets + comparison image set')


# --------------------------------------------------------------------------
# 10b. Slide 18 - Model Architecture (detail + pipeline diagram, no arrows)
# --------------------------------------------------------------------------
s18 = slide_by_title('Model Architecture')
tb18 = shape_by_name(s18, 'TextBox 10')
tb18.left, tb18.top = Inches(0.29), Inches(1.52)
tb18.width, tb18.height = Inches(4.85), Inches(5.25)
S18 = [
    "Six-stage sequential pipeline; stages 3 to 6 run once per clause.",
    ("Ingestion — ", "rule-based text normalisation in the browser via pdf.js; "
     "server-side PyMuPDF text extraction with a Tesseract OCR fallback for "
     "scanned PDFs."),
    ("Segmentation — ", "hybrid regex for structural markers (1., 1.1, (a), (i), "
     "Section / Article) with a spaCy sentenciser fallback; character offsets "
     "are preserved."),
    ("Classifier — ", "fine-tuned nlpaueb/legal-bert-base-uncased "
     "(freak3123/legal-bert-cuad-v1) — four-epoch T4 training on CUAD + a "
     "227-example seed set with class-balanced cross-entropy loss; a TF-IDF + "
     "LogReg artefact and a keyword heuristic are kept as automatic fallbacks."),
    ("NER — ", "hybrid: restrictive regex for the six entity types plus spaCy "
     "en_core_web_sm for DATE / AMOUNT / PARTY recall; regex wins on overlap."),
    ("Rewriter — ", "fine-tuned google/flan-t5-small (freak3123/flan-t5-simplify-v1) "
     "trained on the Plain English Contracts corpus; a template-based simplifier "
     "remains as a fallback."),
    ("Risk flagging — ", "a three-layer rule engine (weighted keyword triggers "
     "and structural regex) sums a score, caps it at 1.0 and maps it to a low, "
     "medium or high level."),
]
fill(tb18.text_frame, bullets(S18, SZ_COL, space_after=6))
drop_pictures(s18)
add_picture(s18, 'model_architecture.png', left=5.32, top=1.66, width=7.75)
print('10b. slide 18 model architecture (no arrows) + pipeline diagram')


# --------------------------------------------------------------------------
# 11. Slide 19 - hardware specs -> Lenovo IdeaPad Slim 3i
# --------------------------------------------------------------------------
s19 = slide_by_title('Software Tools')
tbl19 = next(sh.table for sh in s19.shapes if sh.has_table)
HW = ('Lenovo IdeaPad Slim 3i — Intel Core i7 (13th Gen), 16 GB RAM, '
      '512 GB SSD, Intel Iris Xe integrated graphics; no dedicated GPU')
for r in range(len(tbl19.rows)):
    if 'hardware' in tbl19.cell(r, 0).text.lower():
        cell = tbl19.cell(r, 1)
        size = None
        if cell.text_frame.paragraphs[0].runs:
            size = cell.text_frame.paragraphs[0].runs[0].font.size
        cell.text = HW
        if size and cell.text_frame.paragraphs[0].runs:
            cell.text_frame.paragraphs[0].runs[0].font.size = size
        break
print('11. slide 19 hardware -> Lenovo IdeaPad Slim 3i')


# --------------------------------------------------------------------------
# 12. Slide 22 - Discussion (macro-F1 explained in words, no arrows)
# --------------------------------------------------------------------------
s22 = body_placeholder(slide_by_title('Discussion'))
DISC = [
    "Fine-tuned Legal-BERT attains macro-F1 0.902 across all twelve clause "
    "categories on a held-out 419-clause test set — exceeding the project "
    "macro-F1 target of 0.80.",
    "Lift over baseline is dramatic: TF-IDF + LogReg macro-F1 of 0.61 versus "
    "Legal-BERT macro-F1 of 0.902 (a 48% relative improvement on the same split).",
    "Weakest classes are WARRANTY (F1 0.735) and DISPUTE_RESOLUTION (F1 0.750) — "
    "both show recall = 1.000 with lower precision, i.e. the classifier "
    "over-predicts them; further seed examples are the obvious mitigation.",
    "Rewriter ROUGE climbs monotonically across ten epochs to ROUGE-1 = 0.394, "
    "ROUGE-2 = 0.231 and ROUGE-L = 0.352 — a healthy learning curve on a small "
    "domain corpus.",
    "Engineering wrinkles, openly noted: pdf.js 5.x broke the Next.js build "
    "(pinned to 4.x); Kaggle's current PyTorch lacks kernels for the P100 "
    "(trained on T4 x2 instead).",
    "Risk flagging is kept fully rule-based (no machine learning) so every "
    "flagged clause stays transparent and explainable.",
]
fill(s22.text_frame, [dict(text=t, size=SZ_STD, space_after=8) for t in DISC])
print('12. slide 22 discussion set')


# --------------------------------------------------------------------------
# 12b. Results slide - re-fit the chart so it clears the footer
# --------------------------------------------------------------------------
res = slide_by_title('results')
drop_pictures(res)
_rw = 11.7
add_picture(res, 'results_chart.png', left=(13.333 - _rw) / 2, top=1.74, width=_rw)
print('12b. results chart re-fitted above footer')


# --------------------------------------------------------------------------
# 13. INSERT evaluation-metrics slide before Results
# --------------------------------------------------------------------------
results_idx = next(i for i, s in enumerate(prs.slides)
                   if (s.shapes.title and s.shapes.title.text.strip().lower() == 'results'))
metrics = clone_slide_after(slide_by_title('Discussion'), results_idx - 1)
metrics.shapes.title.text = 'Evaluation Metrics & Formulae'
drop_pictures(metrics)
body_placeholder(metrics).text_frame.clear()
add_picture(metrics, 'metrics_formulae.png', left=1.42, top=1.33, width=10.50)
print('13. evaluation-metrics slide inserted before Results')


# --------------------------------------------------------------------------
# 14. References - replace 8 stubs with all 35, across 3 slides
# --------------------------------------------------------------------------
REFS = [
    "F. Ariai and G. Demartini, “Natural Language Processing for the Legal Domain: "
    "A Survey of Tasks, Datasets, Models and Challenges,” ACM Computing Surveys, 2024 "
    "(arXiv:2410.21306).",
    "“A Survey of Classification Tasks and Approaches for Legal Contracts,” "
    "Artificial Intelligence Review, Springer, 2025.",
    "“A Comprehensive Survey on Legal Summarization: Challenges and Future "
    "Directions,” arXiv:2501.17830, 2025.",
    "D. Ganguly et al., “Legal Text Processing Challenges: NER and Sentence "
    "Boundary Detection,” 2023.",
    "J. Devlin, M.-W. Chang, K. Lee and K. Toutanova, “BERT: Pre-training of Deep "
    "Bidirectional Transformers for Language Understanding,” NAACL-HLT, 2019.",
    "I. Chalkidis, M. Fergadiotis, P. Malakasiotis, N. Aletras and I. Androutsopoulos, "
    "“LEGAL-BERT: The Muppets Straight Out of Law School,” Findings of EMNLP, 2020.",
    "V. Sanh, L. Debut, J. Chaumond and T. Wolf, “DistilBERT, a Distilled Version of "
    "BERT: Smaller, Faster, Cheaper and Lighter,” arXiv:1910.01108, 2019.",
    "C. Raffel, N. Shazeer, A. Roberts et al., “Exploring the Limits of Transfer "
    "Learning with a Unified Text-to-Text Transformer,” JMLR, vol. 21, 2020.",
    "H. W. Chung, L. Hou, S. Longpre et al., “Scaling Instruction-Finetuned "
    "Language Models,” arXiv:2210.11416, 2022.",
    "M. Lewis, Y. Liu, N. Goyal et al., “BART: Denoising Sequence-to-Sequence "
    "Pre-training for Natural Language Generation, Translation and Comprehension,” "
    "ACL, 2020.",
    "D. Hendrycks, C. Burns, A. Chen and S. Ball, “CUAD: An Expert-Annotated NLP "
    "Dataset for Legal Contract Review,” NeurIPS, 2021.",
    "“Cost-Benefit Analysis of Deploying Shallow, Deep Learning and Generative "
    "Models for Legal Text Classification,” ResearchGate / Springer, 2024.",
    "Y. Koreeda and C. D. Manning, “ContractNLI: A Dataset for Document-Level "
    "Natural Language Inference for Contracts,” Findings of EMNLP, 2021.",
    "S. H. Wang et al., “MAUD: An Expert-Annotated Legal NLP Dataset for Merger "
    "Agreement Understanding,” arXiv:2301.00876, 2023.",
    "“ACORD: An Expert-Annotated Retrieval Dataset for Legal Contract Clause "
    "Drafting,” arXiv:2501.06582, 2025.",
    "“Natural Language Processing for Legal Document Review: Categorising Deontic "
    "Modalities in Contracts,” Artificial Intelligence and Law, Springer, 2023.",
    "M. Cemri et al., “Unsupervised Simplification of Legal Texts,” "
    "arXiv:2209.00557, 2022.",
    "“Text Simplification System for Legal Contract Review,” 2024.",
    "“Legal Lay Summarization: Exploring Methods and Data Generation with Large "
    "Language Models,” Artificial Intelligence Review, Springer, 2025.",
    "L. Manor and J. J. Li, “Plain English Summarization of Contracts,” Natural "
    "Legal Language Processing Workshop, 2019.",
    "M. Lippi, P. Pałka, G. Contissa et al., “CLAUDETTE: An Automated Detector of "
    "Potentially Unfair Clauses in Online Terms of Service,” Artificial "
    "Intelligence and Law, vol. 27, 2019.",
    "Ruas et al., “Detecting and Explaining Unfairness in Consumer Contracts "
    "through Memory Networks,” Artificial Intelligence and Law, Springer, 2021.",
    "“A Support System for the Detection of Abusive Clauses in B2C Contracts,” "
    "Artificial Intelligence and Law, Springer, 2024.",
    "“Towards Mitigating Perceived Unfairness in Contracts from a Non-Legal "
    "Perspective,” ACL NLLP Workshop, 2023.",
    "“Software Engineering Meets Legal Texts: LLMs for Auto Detection of Contract "
    "Smells,” ScienceDirect, 2025.",
    "T. Wolf, L. Debut, V. Sanh et al., “Transformers: State-of-the-Art Natural "
    "Language Processing,” EMNLP: System Demonstrations, 2020.",
    "M. Honnibal and I. Montani, “spaCy 2: Natural Language Understanding with "
    "Bloom Embeddings, CNNs and Incremental Parsing,” 2017.",
    "R. Smith, “An Overview of the Tesseract OCR Engine,” Int. Conf. on Document "
    "Analysis and Recognition (ICDAR), 2007.",
    "I. Loshchilov and F. Hutter, “Decoupled Weight Decay Regularization,” "
    "ICLR, 2019.",
    "K. Papineni, S. Roukos, T. Ward and W.-J. Zhu, “BLEU: A Method for Automatic "
    "Evaluation of Machine Translation,” ACL, 2002.",
    "C.-Y. Lin, “ROUGE: A Package for Automatic Evaluation of Summaries,” Text "
    "Summarization Branches Out, 2004.",
    "M. Abadi et al., “TensorFlow: Large-Scale Machine Learning on Heterogeneous "
    "Systems,” 2015.",
    "I. Chalkidis et al., “LexGLUE: A Benchmark Dataset for Legal Language "
    "Understanding in English,” arXiv:2110.00976, 2021.",
    "“LEDGAR: A Large-Scale Multi-Label Corpus for Text Classification of Legal "
    "Provisions in Contracts,” LREC, 2020.",
    "K. Vuthoo, S. Khetarpaul and S. Mishra, “Efficient Clause Identification in "
    "Contracts Using NLP and Web-Sourced Data,” ICDAM 2025, Springer LNNS "
    "vol. 1594, 2026.",
]


def fill_refs(slide, start, end, title):
    slide.shapes.title.text = title
    specs = []
    for i in range(start, end):
        specs.append(dict(size=SZ_REF, space_after=4, bullet=False, hanging=True,
                           runs=[(f'[{i + 1}]  ', True, NAVY), (REFS[i], False)]))
    fill(body_placeholder(slide).text_frame, specs)


ref_slide = slide_by_title('References')
ref_idx = list(prs.slides).index(ref_slide)
ref2 = clone_slide_after(ref_slide, ref_idx)
ref3 = clone_slide_after(ref_slide, ref_idx + 1)
fill_refs(ref_slide, 0, 12, 'References')
fill_refs(ref2, 12, 24, 'References (cont.)')
fill_refs(ref3, 24, 35, 'References (cont.)')
print('14. references expanded to all 35 across 3 slides')


# --------------------------------------------------------------------------
# 15. renumber slide-number placeholders and save
# --------------------------------------------------------------------------
renumber()
prs.save(OUT)
print(f'\nSaved {OUT}  -  {len(prs.slides)} slides total')
